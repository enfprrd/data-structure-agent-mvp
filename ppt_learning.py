from __future__ import annotations

import json
import re
import tempfile
from dataclasses import asdict, dataclass, field
from io import BytesIO
from pathlib import Path
from typing import Any

from PIL import Image, ImageDraw, ImageFont
from conversation_context import format_message_history
from llm import DeepSeekClient, DeepSeekError
from rag import MarkdownKeywordRetriever
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN


TOKEN_PATTERN = re.compile(r"[A-Za-z_][A-Za-z0-9_]*|[\u4e00-\u9fff]+")
CODE_BLOCK_PATTERN = re.compile(r"```[\s\S]*?```|(?:^|\n)(?: {4}|\t).+", re.MULTILINE)


@dataclass
class SlideCard:
    deck_id: str
    slide_id: int
    title: str
    raw_text: str
    notes: str
    tables: list[list[list[str]]] = field(default_factory=list)
    code_blocks: list[str] = field(default_factory=list)
    summary: str = ""
    keywords: list[str] = field(default_factory=list)
    concept_type: str = "unknown"

    @property
    def has_text(self) -> bool:
        return bool(self.raw_text.strip() or self.notes.strip() or self.tables)

    def to_dict(self) -> dict[str, Any]:
        return asdict(self)


@dataclass
class ContextPack:
    current_slide: SlideCard
    prev_slide: dict[str, Any] | None
    next_slide: dict[str, Any] | None
    deck_outline: list[dict[str, Any]]
    retrieved_slides: list[SlideCard]
    conversation_summary: str
    textbook_contexts: list[dict[str, str]] = field(default_factory=list)


def parse_pptx_to_slide_cards(
    deck_id: str,
    pptx_bytes: bytes,
    client: DeepSeekClient | None = None,
) -> list[SlideCard]:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise RuntimeError("缺少 python-pptx，请先安装 requirements.txt 中的依赖。") from exc

    presentation = Presentation(BytesIO(pptx_bytes))
    cards: list[SlideCard] = []
    for index, slide in enumerate(presentation.slides, start=1):
        card = _extract_slide(deck_id, index, slide)
        if card.has_text:
            _enrich_slide_card(card, client)
        else:
            card.summary = "该页没有可提取的文本内容，第一阶段不猜测图片信息。"
            card.keywords = ["text_missing"]
            card.concept_type = "text_missing"
        cards.append(card)
    return cards


def render_pptx_slide_images(
    deck_id: str,
    pptx_bytes: bytes,
    output_root: Path | str,
    width: int = 1440,
    height: int = 810,
) -> list[str]:
    """Render PPTX slides to local PNG files for user-facing preview.

    This is display-only. The generated images are never sent to DeepSeek.
    """
    output_dir = (Path(output_root) / deck_id).resolve()
    output_dir.mkdir(parents=True, exist_ok=True)

    existing = sorted(output_dir.glob("slide_*.png"))
    if existing:
        return [str(path) for path in existing]

    rendered = _render_pptx_slide_images_with_powerpoint(deck_id, pptx_bytes, output_dir, width, height)
    if rendered:
        return rendered
    return _render_pptx_slide_images_with_pillow(deck_id, pptx_bytes, output_dir, width, height)


def _render_pptx_slide_images_with_powerpoint(
    deck_id: str,
    pptx_bytes: bytes,
    output_dir: Path,
    width: int,
    height: int,
) -> list[str]:
    try:
        import pythoncom
        import win32com.client
    except ImportError:
        return []

    with tempfile.TemporaryDirectory(prefix="ppt_render_") as temp_dir:
        temp_path = Path(temp_dir) / f"{deck_id}.pptx"
        temp_path.write_bytes(pptx_bytes)

        powerpoint = None
        presentation = None
        pythoncom.CoInitialize()
        try:
            powerpoint = win32com.client.DispatchEx("PowerPoint.Application")
            presentation = powerpoint.Presentations.Open(
                str(temp_path),
                ReadOnly=True,
                Untitled=False,
                WithWindow=False,
            )
            rendered: list[str] = []
            for index in range(1, presentation.Slides.Count + 1):
                image_path = output_dir / f"slide_{index:03d}.png"
                presentation.Slides(index).Export(str(image_path), "PNG", width, height)
                rendered.append(str(image_path))
            return rendered
        except Exception:
            return []
        finally:
            if presentation is not None:
                try:
                    presentation.Close()
                except Exception:
                    pass
            if powerpoint is not None:
                try:
                    powerpoint.Quit()
                except Exception:
                    pass
            pythoncom.CoUninitialize()


def _render_pptx_slide_images_with_pillow(
    deck_id: str,
    pptx_bytes: bytes,
    output_dir: Path,
    width: int,
    height: int,
) -> list[str]:
    try:
        from pptx import Presentation
    except ImportError:
        return []

    try:
        presentation = Presentation(BytesIO(pptx_bytes))
    except Exception:
        return []

    slide_width_px, slide_height_px = width, height
    slide_width_emu = presentation.slide_width
    slide_height_emu = presentation.slide_height

    rendered: list[str] = []
    for index, slide in enumerate(presentation.slides, start=1):
        image = Image.new("RGB", (slide_width_px, slide_height_px), "white")
        draw = ImageDraw.Draw(image)
        _draw_ppt_slide_background(draw, slide, slide_width_px, slide_height_px)
        for shape in slide.shapes:
            _draw_ppt_shape(draw, image, shape, slide_width_px, slide_height_px, slide_width_emu, slide_height_emu)
        image_path = output_dir / f"slide_{index:03d}.png"
        image.save(image_path)
        rendered.append(str(image_path))
    return rendered


def discover_local_pptx_files(directory: Path | str) -> list[Path]:
    base = Path(directory)
    if not base.exists() or not base.is_dir():
        return []
    return sorted(
        [path for path in base.iterdir() if path.is_file() and path.suffix.lower() == ".pptx"],
        key=lambda path: path.name.lower(),
    )


def _draw_ppt_slide_background(draw: ImageDraw.ImageDraw, slide: Any, width: int, height: int) -> None:
    try:
        background = slide.background
        fill = background.fill
        if getattr(fill, "type", None) and getattr(fill, "fore_color", None) is not None:
            color = getattr(fill.fore_color, "rgb", None)
            if color is not None:
                draw.rectangle((0, 0, width, height), fill=f"#{color}")
    except Exception:
        return


def _draw_ppt_shape(
    draw: ImageDraw.ImageDraw,
    image: Image.Image,
    shape: Any,
    width: int,
    height: int,
    slide_width_emu: int,
    slide_height_emu: int,
) -> None:
    if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PLACEHOLDER:
        if getattr(shape, "has_text_frame", False):
            _draw_ppt_text_frame(draw, shape, width, height, slide_width_emu, slide_height_emu)
        return

    if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE:
        _draw_ppt_picture(image, shape, width, height, slide_width_emu, slide_height_emu)
        return

    if getattr(shape, "has_table", False):
        _draw_ppt_table(draw, shape.table, width, height, slide_width_emu, slide_height_emu)
        return

    if getattr(shape, "has_text_frame", False):
        _draw_ppt_text_frame(draw, shape, width, height, slide_width_emu, slide_height_emu)
        return

    # Basic fallback for simple auto-shapes.
    try:
        left, top, right, bottom = _shape_box(shape, width, height, slide_width_emu, slide_height_emu)
    except Exception:
        return
    try:
        fill = getattr(shape.fill, "fore_color", None)
        fill_color = f"#{fill.rgb}" if fill is not None and getattr(fill, "rgb", None) is not None else None
    except Exception:
        fill_color = None
    try:
        line = getattr(shape.line, "color", None)
        line_color = f"#{line.rgb}" if line is not None and getattr(line, "rgb", None) is not None else "#d1d5db"
    except Exception:
        line_color = "#d1d5db"
    draw.rectangle((left, top, right, bottom), fill=fill_color, outline=line_color, width=2)


def _draw_ppt_picture(
    image: Image.Image,
    shape: Any,
    width: int,
    height: int,
    slide_width_emu: int,
    slide_height_emu: int,
) -> None:
    try:
        from PIL import Image as PILImage
    except Exception:
        return

    try:
        left, top, right, bottom = _shape_box(shape, width, height, slide_width_emu, slide_height_emu)
        blob = shape.image.blob
        picture = PILImage.open(BytesIO(blob)).convert("RGBA")
        box_width = max(1, right - left)
        box_height = max(1, bottom - top)
        picture.thumbnail((box_width, box_height), PILImage.Resampling.LANCZOS)
        offset_x = left + max(0, (box_width - picture.width) // 2)
        offset_y = top + max(0, (box_height - picture.height) // 2)
        image.paste(picture, (offset_x, offset_y), picture)
    except Exception:
        return


def _draw_ppt_table(
    draw: ImageDraw.ImageDraw,
    table: Any,
    width: int,
    height: int,
    slide_width_emu: int,
    slide_height_emu: int,
) -> None:
    try:
        rows = list(table.rows)
        cols = list(table.columns)
        if not rows or not cols:
            return
        left, top, right, bottom = _table_box(table, width, height, slide_width_emu, slide_height_emu)
        box_width = max(1, right - left)
        box_height = max(1, bottom - top)
        row_heights = [box_height / len(rows)] * len(rows)
        col_widths = [box_width / len(cols)] * len(cols)
        y = top
        for row_index, row in enumerate(rows):
            x = left
            row_height = row_heights[row_index]
            for col_index, cell in enumerate(row.cells):
                cell_width = col_widths[col_index]
                cell_box = (x, y, x + cell_width, y + row_height)
                draw.rectangle(cell_box, fill="white", outline="#d1d5db", width=1)
                _draw_cell_text(draw, cell.text, cell_box, default_size=16, align="left")
                x += cell_width
            y += row_height
    except Exception:
        return


def _draw_ppt_text_frame(
    draw: ImageDraw.ImageDraw,
    shape: Any,
    width: int,
    height: int,
    slide_width_emu: int,
    slide_height_emu: int,
) -> None:
    try:
        left, top, right, bottom = _shape_box(shape, width, height, slide_width_emu, slide_height_emu)
        text_frame = shape.text_frame
        paragraphs = list(text_frame.paragraphs)
        if not paragraphs:
            return
        padding_left = _emu_to_px(getattr(text_frame, "margin_left", 0), slide_width_emu, width)
        padding_right = _emu_to_px(getattr(text_frame, "margin_right", 0), slide_width_emu, width)
        padding_top = _emu_to_px(getattr(text_frame, "margin_top", 0), slide_height_emu, height)
        padding_bottom = _emu_to_px(getattr(text_frame, "margin_bottom", 0), slide_height_emu, height)
        inner = (
            left + padding_left,
            top + padding_top,
            max(left + padding_left + 1, right - padding_right),
            max(top + padding_top + 1, bottom - padding_bottom),
        )
        lines: list[tuple[str, dict[str, Any]]] = []
        for paragraph in paragraphs:
            text = "".join(run.text for run in paragraph.runs).strip() or paragraph.text.strip()
            if not text:
                continue
            line_meta = {
                "size": _paragraph_font_size(paragraph),
                "bold": _paragraph_is_bold(paragraph),
                "italic": _paragraph_is_italic(paragraph),
                "color": _paragraph_color(paragraph),
                "align": paragraph.alignment or PP_ALIGN.LEFT,
            }
            lines.append((text, line_meta))
        if not lines:
            return
        _draw_wrapped_text_block(draw, lines, inner)
    except Exception:
        return


def _draw_cell_text(
    draw: ImageDraw.ImageDraw,
    text: str,
    box: tuple[float, float, float, float],
    default_size: int = 16,
    align: str = "center",
) -> None:
    lines = [(line, {"size": default_size, "bold": False, "italic": False, "color": "#111827", "align": PP_ALIGN.LEFT}) for line in text.splitlines() if line.strip()]
    if not lines and text.strip():
        lines = [(text.strip(), {"size": default_size, "bold": False, "italic": False, "color": "#111827", "align": PP_ALIGN.LEFT})]
    if lines:
        _draw_wrapped_text_block(draw, lines, box, force_align=align)


def _draw_wrapped_text_block(
    draw: ImageDraw.ImageDraw,
    lines: list[tuple[str, dict[str, Any]]],
    box: tuple[float, float, float, float],
    force_align: str | None = None,
) -> None:
    left, top, right, bottom = map(int, box)
    max_width = max(1, right - left)
    max_height = max(1, bottom - top)
    fitted: list[tuple[str, ImageFont.FreeTypeFont, str, str]] = []
    for text, meta in lines:
        font_size = max(12, int(meta.get("size") or 16))
        font = _load_ppt_font(font_size, bool(meta.get("bold")), bool(meta.get("italic")))
        wrapped = _wrap_text(draw, text, font, max_width)
        if not wrapped:
            wrapped = [text]
        align = force_align or _alignment_to_name(meta.get("align"))
        color = str(meta.get("color") or "#111827")
        fitted.append(("\n".join(wrapped), font, align, color))

    total_height = sum(_measure_multiline(draw, text, font, max_width)[1] + 6 for text, font, _, _ in fitted)
    scale = min(1.0, max_height / max(total_height, 1))
    if scale < 1.0:
        resized: list[tuple[str, ImageFont.FreeTypeFont, str, str]] = []
        for text, font, align, color in fitted:
            target_size = max(12, int(round(_font_size_px(font) * scale)))
            resized.append((text, _load_ppt_font(target_size, font.path.endswith("b.ttf") if getattr(font, "path", "") else False, False), align, color))
        fitted = resized

    cursor_y = top
    for text, font, align, color in fitted:
        text_width, text_height = _measure_multiline(draw, text, font, max_width)
        if align == "center":
            x = left + max(0, (max_width - text_width) // 2)
        elif align == "right":
            x = left + max(0, max_width - text_width)
        else:
            x = left
        y = cursor_y
        draw.multiline_text((x, y), text, fill=color, font=font, spacing=4)
        cursor_y += text_height + 6


def _shape_box(
    shape: Any,
    width: int,
    height: int,
    slide_width_emu: int,
    slide_height_emu: int,
) -> tuple[int, int, int, int]:
    return (
        _emu_to_px(getattr(shape, "left", 0), slide_width_emu, width),
        _emu_to_px(getattr(shape, "top", 0), slide_height_emu, height),
        _emu_to_px(getattr(shape, "left", 0) + getattr(shape, "width", 0), slide_width_emu, width),
        _emu_to_px(getattr(shape, "top", 0) + getattr(shape, "height", 0), slide_height_emu, height),
    )


def _table_box(
    table: Any,
    width: int,
    height: int,
    slide_width_emu: int,
    slide_height_emu: int,
) -> tuple[int, int, int, int]:
    parent = getattr(table, "_parent", None)
    if parent is not None:
        return _shape_box(parent, width, height, slide_width_emu, slide_height_emu)
    return (0, 0, width, height)


def _emu_to_px(value: Any, axis_total_emu: int, axis_pixels: int) -> int:
    try:
        if axis_total_emu == 0:
            return 0
        return int(round(float(value) / float(axis_total_emu) * axis_pixels))
    except Exception:
        return 0


def _paragraph_font_size(paragraph: Any, default: int = 20) -> int:
    for run in paragraph.runs:
        size = getattr(getattr(run, "font", None), "size", None)
        if size is not None:
            try:
                return max(10, int(round(size.pt * 96 / 72)))
            except Exception:
                continue
    return default


def _paragraph_is_bold(paragraph: Any) -> bool:
    for run in paragraph.runs:
        if getattr(getattr(run, "font", None), "bold", None):
            return True
    return False


def _paragraph_is_italic(paragraph: Any) -> bool:
    for run in paragraph.runs:
        if getattr(getattr(run, "font", None), "italic", None):
            return True
    return False


def _paragraph_color(paragraph: Any) -> str:
    for run in paragraph.runs:
        font = getattr(run, "font", None)
        color = getattr(font, "color", None)
        rgb = getattr(color, "rgb", None)
        if rgb is not None:
            return f"#{rgb}"
    return "#111827"


def _alignment_to_name(value: Any) -> str:
    if value in {PP_ALIGN.CENTER, "center"}:
        return "center"
    if value in {PP_ALIGN.RIGHT, "right"}:
        return "right"
    return "left"


def _load_ppt_font(size_px: int, bold: bool = False, italic: bool = False) -> ImageFont.FreeTypeFont:
    candidates = [
        (r"C:\\Windows\\Fonts\\msyhbd.ttc" if bold else r"C:\\Windows\\Fonts\\msyh.ttc"),
        (r"C:\\Windows\\Fonts\\simhei.ttf" if bold else r"C:\\Windows\\Fonts\\simhei.ttf"),
        (r"C:\\Windows\\Fonts\\Noto Sans SC Bold (TrueType).otf" if bold else r"C:\\Windows\\Fonts\\Noto Sans SC (TrueType).otf"),
        (r"C:\\Windows\\Fonts\\arialbd.ttf" if bold else r"C:\\Windows\\Fonts\\arial.ttf"),
    ]
    for path in candidates:
        font_path = Path(path)
        if font_path.exists():
            try:
                return ImageFont.truetype(str(font_path), size=max(12, size_px))
            except Exception:
                continue
    return ImageFont.load_default()


def _font_size_px(font: ImageFont.FreeTypeFont) -> int:
    try:
        return int(round(font.size))
    except Exception:
        return 16


def _measure_multiline(draw: ImageDraw.ImageDraw, text: str, font: ImageFont.FreeTypeFont, max_width: int) -> tuple[int, int]:
    wrapped = _wrap_text(draw, text, font, max_width) or [text]
    bbox = draw.multiline_textbbox((0, 0), "\n".join(wrapped), font=font, spacing=4)
    return max(1, int(bbox[2] - bbox[0])), max(1, int(bbox[3] - bbox[1]))


def _wrap_text(draw: ImageDraw.ImageDraw, text: str, font: ImageFont.FreeTypeFont, max_width: int) -> list[str]:
    paragraphs = text.splitlines() or [text]
    wrapped: list[str] = []
    for paragraph in paragraphs:
        stripped = paragraph.strip()
        if not stripped:
            wrapped.append("")
            continue
        tokens = stripped.split()
        if len(tokens) == 1:
            tokens = list(stripped)
        current = ""
        for token in tokens:
            candidate = token if not current else f"{current} {token}" if " " in stripped else f"{current}{token}"
            if _text_width(draw, candidate, font) <= max_width or not current:
                current = candidate
            else:
                wrapped.append(current)
                current = token
        if current:
            wrapped.append(current)
    return wrapped


def _text_width(draw: ImageDraw.ImageDraw, text: str, font: ImageFont.FreeTypeFont) -> int:
    bbox = draw.textbbox((0, 0), text, font=font)
    return max(1, int(bbox[2] - bbox[0]))


def build_context_pack(
    slide_cards: list[SlideCard],
    current_slide_id: int,
    question: str,
    conversation_messages: list[dict[str, str]] | None = None,
    textbook_retriever: MarkdownKeywordRetriever | None = None,
    top_k: int = 3,
) -> ContextPack:
    if not slide_cards:
        raise ValueError("slide_cards 不能为空")

    index = min(max(current_slide_id - 1, 0), len(slide_cards) - 1)
    current_slide = slide_cards[index]
    retrieved = retrieve_related_slides(slide_cards, question, current_slide.slide_id, top_k=top_k)
    textbook_contexts = (
        textbook_retriever.retrieve(question, top_k=2, boost_terms=current_slide.keywords)
        if textbook_retriever is not None and question.strip()
        else []
    )
    return ContextPack(
        current_slide=current_slide,
        prev_slide=_slide_brief(slide_cards[index - 1]) if index > 0 else None,
        next_slide=_slide_brief(slide_cards[index + 1]) if index < len(slide_cards) - 1 else None,
        deck_outline=[{"slide_id": card.slide_id, "title": card.title or f"第 {card.slide_id} 页"} for card in slide_cards],
        retrieved_slides=retrieved,
        conversation_summary=summarize_conversation(conversation_messages or []),
        textbook_contexts=textbook_contexts,
    )


def answer_with_context_pack(
    client: DeepSeekClient,
    context_pack: ContextPack,
    question: str,
) -> str:
    messages = [
        {
            "role": "system",
            "content": (
                "你是 PPT-Aware 数据结构学习助教。你只能读取用户提供的 PPT 文本、备注、表格、"
                "检索到的相关页和本地教材片段；不能声称看过 PPT 图片或截图。"
            ),
        },
        {"role": "user", "content": build_ppt_user_prompt(context_pack, question)},
    ]
    return client.chat(messages, temperature=0.2, max_tokens=2200)


def build_ppt_user_prompt(context_pack: ContextPack, question: str) -> str:
    pack = context_pack_to_prompt_dict(context_pack)
    return f"""
用户当前正在阅读 PPT 第 {context_pack.current_slide.slide_id} 页。

用户问题：
{question}

ContextPack（只包含当前页完整文本、相邻页摘要、目录、相关页和少量教材片段）：
{json.dumps(pack, ensure_ascii=False, indent=2)}

回答规则：
1. 优先依据 current_slide。
2. 其次依据 prev_slide、next_slide 和 retrieved_slides。
3. 再参考 textbook_contexts。
4. 如果 PPT 文本没有相关信息，明确说“不在当前 PPT 文本中”，不要猜图片内容。
5. 每个关键结论必须标注来源页码，例如“来源：第 3 页”或“来源：第 3、5 页”。
6. 如果 current_slide.concept_type 是 text_missing，先提示该页可提取文本不足。
7. 不要说你看到了图片、图表截图或视觉布局；第一阶段只做文本版。
"""


def context_pack_to_prompt_dict(context_pack: ContextPack) -> dict[str, Any]:
    return {
        "current_slide": context_pack.current_slide.to_dict(),
        "prev_slide": context_pack.prev_slide,
        "next_slide": context_pack.next_slide,
        "deck_outline": context_pack.deck_outline,
        "retrieved_slides": [
            {
                "slide_id": card.slide_id,
                "title": card.title,
                "summary": card.summary,
                "keywords": card.keywords,
                "raw_text_excerpt": _shorten(card.raw_text, 900),
            }
            for card in context_pack.retrieved_slides
        ],
        "conversation_summary": context_pack.conversation_summary,
        "textbook_contexts": context_pack.textbook_contexts,
    }


def retrieve_related_slides(
    slide_cards: list[SlideCard],
    question: str,
    current_slide_id: int,
    top_k: int = 3,
) -> list[SlideCard]:
    query_tokens = _tokenize(question)
    if not query_tokens:
        query_tokens = _tokenize(" ".join(_find_slide(slide_cards, current_slide_id).keywords))
    scored: list[tuple[int, SlideCard]] = []
    for card in slide_cards:
        if card.slide_id == current_slide_id:
            continue
        haystack = " ".join(
            [
                card.title,
                card.raw_text,
                card.notes,
                card.summary,
                " ".join(card.keywords),
                _tables_to_text(card.tables),
            ]
        ).lower()
        score = sum(haystack.count(token.lower()) for token in query_tokens)
        for keyword in card.keywords:
            if keyword and keyword.lower() in question.lower():
                score += 4
        if score > 0:
            scored.append((score, card))
    scored.sort(key=lambda item: (-item[0], item[1].slide_id))
    return [card for _, card in scored[:top_k]]


def summarize_conversation(messages: list[dict[str, object]], max_items: int = 6) -> str:
    return format_message_history(
        messages,
        include_code_blocks=False,
        max_messages=max_items,
        per_message_limit=180,
        empty_text="暂无会话上下文。",
        separator="\n",
    )


def _extract_slide(deck_id: str, slide_id: int, slide: Any) -> SlideCard:
    title = _extract_title(slide)
    text_parts: list[str] = []
    tables: list[list[list[str]]] = []

    for shape in slide.shapes:
        if getattr(shape, "has_table", False):
            tables.append(_extract_table(shape.table))
            continue
        if getattr(shape, "has_text_frame", False):
            text = shape.text_frame.text.strip()
            if text:
                text_parts.append(text)

    notes = _extract_notes(slide)
    raw_text = _dedupe_join([title, *text_parts])
    code_blocks = [match.group(0).strip() for match in CODE_BLOCK_PATTERN.finditer(raw_text)]
    return SlideCard(
        deck_id=deck_id,
        slide_id=slide_id,
        title=title,
        raw_text=raw_text,
        notes=notes,
        tables=tables,
        code_blocks=code_blocks,
    )


def _extract_title(slide: Any) -> str:
    title_shape = getattr(slide.shapes, "title", None)
    if title_shape is not None and getattr(title_shape, "has_text_frame", False):
        return title_shape.text_frame.text.strip()
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False):
            text = shape.text_frame.text.strip()
            if text:
                return text.splitlines()[0].strip()
    return ""


def _extract_table(table: Any) -> list[list[str]]:
    rows: list[list[str]] = []
    for row in table.rows:
        rows.append([cell.text.strip() for cell in row.cells])
    return rows


def _extract_notes(slide: Any) -> str:
    try:
        notes_slide = slide.notes_slide
        text_frame = notes_slide.notes_text_frame
    except Exception:
        return ""
    return text_frame.text.strip() if text_frame is not None else ""


def _enrich_slide_card(card: SlideCard, client: DeepSeekClient | None) -> None:
    fallback_summary = _extractive_summary(card)
    card.summary = fallback_summary
    card.keywords = _extract_keywords(card)
    card.concept_type = "text"

    if client is None:
        return
    try:
        raw = client.chat(
            [
                {
                    "role": "system",
                    "content": "你只输出 JSON，不要解释。字段：summary, keywords, concept_type。",
                },
                {
                    "role": "user",
                    "content": (
                        "请基于这一页 PPT 的可提取文本生成元数据，不要猜测图片内容。\n"
                        f"标题：{card.title}\n"
                        f"正文：{_shorten(card.raw_text, 3500)}\n"
                        f"备注：{_shorten(card.notes, 1200)}\n"
                        f"表格：{_shorten(_tables_to_text(card.tables), 1200)}\n"
                        '输出 JSON 示例：{"summary":"...","keywords":["..."],"concept_type":"definition|algorithm|example|exercise|overview|text"}'
                    ),
                },
            ],
            temperature=0,
            max_tokens=500,
        )
    except DeepSeekError:
        return
    meta = _parse_slide_meta(raw)
    if meta:
        card.summary = meta.get("summary") or fallback_summary
        card.keywords = meta.get("keywords") or card.keywords
        card.concept_type = meta.get("concept_type") or card.concept_type


def _parse_slide_meta(raw: str) -> dict[str, Any]:
    cleaned = raw.strip()
    cleaned = re.sub(r"^```(?:json)?", "", cleaned).strip()
    cleaned = re.sub(r"```$", "", cleaned).strip()
    try:
        data = json.loads(cleaned)
    except json.JSONDecodeError:
        return {}
    if not isinstance(data, dict):
        return {}
    keywords = data.get("keywords")
    if not isinstance(keywords, list):
        keywords = []
    return {
        "summary": str(data.get("summary") or "").strip()[:500],
        "keywords": [str(item).strip()[:30] for item in keywords if str(item).strip()][:10],
        "concept_type": str(data.get("concept_type") or "text").strip()[:30],
    }


def _extractive_summary(card: SlideCard) -> str:
    text = " ".join([card.title, card.raw_text, card.notes, _tables_to_text(card.tables)]).strip()
    if not text:
        return "该页没有可提取的文本内容，第一阶段不猜测图片信息。"
    sentences = re.split(r"(?<=[。！？.!?])\s+|\n+", text)
    summary = " ".join(sentence.strip() for sentence in sentences if sentence.strip())[:240]
    return summary or _shorten(text, 240)


def _extract_keywords(card: SlideCard, limit: int = 8) -> list[str]:
    text = " ".join([card.title, card.raw_text, card.notes, _tables_to_text(card.tables)])
    tokens = _tokenize(text)
    counts: dict[str, int] = {}
    for token in tokens:
        if len(token) <= 1 and not token.isascii():
            continue
        counts[token] = counts.get(token, 0) + 1
    return [token for token, _ in sorted(counts.items(), key=lambda item: (-item[1], item[0]))[:limit]]


def _tokenize(text: str) -> list[str]:
    return [token.lower() for token in TOKEN_PATTERN.findall(text) if token.strip()]


def _slide_brief(card: SlideCard) -> dict[str, Any]:
    return {
        "slide_id": card.slide_id,
        "title": card.title,
        "summary": card.summary,
        "keywords": card.keywords,
        "concept_type": card.concept_type,
    }


def _find_slide(slide_cards: list[SlideCard], slide_id: int) -> SlideCard:
    for card in slide_cards:
        if card.slide_id == slide_id:
            return card
    return slide_cards[0]


def _tables_to_text(tables: list[list[list[str]]]) -> str:
    lines: list[str] = []
    for table_index, table in enumerate(tables, start=1):
        lines.append(f"表格 {table_index}:")
        for row in table:
            lines.append(" | ".join(row))
    return "\n".join(lines)


def _dedupe_join(parts: list[str]) -> str:
    lines: list[str] = []
    seen: set[str] = set()
    for part in parts:
        for line in part.splitlines():
            normalized = line.strip()
            if not normalized or normalized in seen:
                continue
            lines.append(normalized)
            seen.add(normalized)
    return "\n".join(lines)


def _shorten(text: str, limit: int) -> str:
    text = text.strip()
    return text if len(text) <= limit else text[:limit].rstrip() + "..."
